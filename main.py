import nltk
nltk.download('punkt')
nltk.download('punkt_tab')
import pdfkit
from flask import Flask, render_template_string, request, send_file
from transformers import T5ForConditionalGeneration, T5Tokenizer, pipeline
from nltk.tokenize import sent_tokenize, word_tokenize
from pyngrok import ngrok
import io
from flask import render_template_string
from gtts import gTTS
from moviepy.editor import *
import os
import requests
from bs4 import BeautifulSoup
import json
from pptx import Presentation


UNSPLASH_ACCESS_KEY = "CRhb-o5AP9jrv3h-QhBMslPHvoCd7chi2_OVRzVdHt0"

app = Flask(__name__)

# ==========================Video generation
def fetch_images_from_unsplash(keywords, num_images=3):
    images = []
    for keyword in keywords:
        response = requests.get(
            f"https://api.unsplash.com/search/photos",
            params={"query": keyword, "per_page": num_images},
            headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
        )
        data = response.json()
        for photo in data['results']:
            images.append(photo['urls']['regular'])
        if len(images) >= num_images:
            break
    return images[:num_images]

def generate_video_from_text_and_images(summary, key_points):
    text_to_speech = f"Summary: {summary}. Key Points: {', '.join(key_points)}."

    tts = gTTS(text=text_to_speech, lang='en')
    audio_file = "output_audio.mp3"
    tts.save(audio_file)

    keywords = summary.split()[:3]
    images = fetch_images_from_unsplash(keywords)

    audio_clip = AudioFileClip(audio_file)
    image_clips = []
    for img_url in images:
        image = ImageClip(img_url).set_duration(audio_clip.duration / len(images))
        image = image.resize(height=720)
        image_clips.append(image)

    video = concatenate_videoclips(image_clips, method="compose").set_audio(audio_clip)

    video_file = "study_material_video.mp4"
    video.write_videofile(video_file, fps=24)

    os.remove(audio_file)

    return video_file
# ==========================Video generation





# SCRAPPING==========================
def scrape_geeksforgeeks(topic):
    search_url = f"https://www.geeksforgeeks.org/{topic.replace(' ', '-').lower()}"
    print(search_url)
    response = requests.get(search_url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, 'html.parser')
        content = ""
        for p in soup.find_all('p'):
            content += p.get_text()
        return content
    else:
        return "Error: Unable to fetch content"

def scrape_wikipedia(topic):
    url = f"https://en.wikipedia.org/wiki/{topic.replace(' ', '_')}"
    print(url)
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, 'html.parser')
        content = ""
        for p in soup.find_all('p'):
            content += p.get_text()
        return content
    else:
        return "Error: Unable to fetch content"

def save_content_to_file(content, filename):
    with open(filename, "w", encoding="utf-8") as file:
        file.write(content)

def scrape_content(topic):
    wikipedia_content = scrape_wikipedia(topic)
    print("Wiki")
    print(wikipedia_content[:500])
    gfg_content = scrape_geeksforgeeks(topic)
    print("GfG")
    print(gfg_content[:500])
    combined_content = gfg_content + "\n\n" + wikipedia_content  # Combine content
    filename = f"{topic.replace(' ', '-').lower()}_content.txt"
    save_content_to_file(combined_content, filename)
    print(f"Content saved to {filename}")
    return combined_content[:2000]

# SCRAPPING==========================







# ================split data
def split_into_chunks(text, tokenizer, max_tokens=512):
    words = word_tokenize(text)
    chunks = []
    current_chunk = []
    current_length = 0

    for word in words:
        current_length += len(tokenizer.tokenize(word))
        if current_length <= max_tokens:
            current_chunk.append(word)
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]
            current_length = len(tokenizer.tokenize(word))

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

def generate_summary(paragraph):
    # summarizer = pipeline("summarization")

    summarizer = pipeline(
        "summarization",
        model="sshleifer/distilbart-cnn-12-6",
        revision="a4f8f3e"
    )

    tokenizer = summarizer.tokenizer
    chunks = split_into_chunks(paragraph, tokenizer)
    summary_text = [summarizer(chunk)[0]['summary_text'] for chunk in chunks]
    return " ".join(summary_text)

def generate_key_points(paragraph):
    tokenizer = T5Tokenizer.from_pretrained('valhalla/t5-small-e2e-qg')
    chunks = split_into_chunks(paragraph, tokenizer)
    key_points = [sent_tokenize(chunk) for chunk in chunks]
    return [sent for chunk in key_points for sent in chunk if len(sent.split()) > 5]

def generate_questions(paragraph):
    model = T5ForConditionalGeneration.from_pretrained('valhalla/t5-small-e2e-qg')
    tokenizer = T5Tokenizer.from_pretrained('valhalla/t5-small-e2e-qg')
    chunks = split_into_chunks(paragraph, tokenizer)
    all_questions = []
    for chunk in chunks:
        input_text = f"generate questions: {chunk}"
        input_ids = tokenizer.encode(input_text, return_tensors="pt", max_length=512, truncation=True)
        outputs = model.generate(input_ids, max_length=150, num_return_sequences=3, num_beams=5)
        questions = [tokenizer.decode(output, skip_special_tokens=True) for output in outputs]
        all_questions.extend(questions)
    return all_questions


def generate_study_material(content):
    return {
        'Summary': generate_summary(content),
        'Key Points': generate_key_points(content),
        'Questions': generate_questions(content)
    }

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        topic = request.form['topic']

        scraped_content = scrape_content(topic)
        print("scraped_content\/\/\/\/\/\/\/\/\/\/\/\/\/\/")
        print(scraped_content)
        study_material = generate_study_material(scraped_content)

        return render_template_string('''

  <!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Study Material Generator</title>
    <style>
        /* General reset */
        body {
            margin: 0;
            padding: 0;
            font-family: 'Arial', sans-serif;
            background: #f8f9fa;
        }

        .container {
            max-width: 800px;
            margin: 50px auto;
            padding: 20px;
            background: #fff;
            border-radius: 10px;
            box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1);
        }

        h1 {
            font-size: 2rem;
            color: #333;
            text-align: center;
            margin-bottom: 20px;
        }

        .form-group label {
            display: block;
            font-size: 1rem;
            font-weight: 600;
            margin-bottom: 5px;
            color: #555;
        }

        .form-group input {
            width: 100%;
            padding: 10px;
            font-size: 1rem;
            border: 1px solid #ccc;
            border-radius: 5px;
        }

        .btn {
            display: block;
            width: 100%;
            padding: 10px 15px;
            font-size: 1rem;
            font-weight: bold;
            text-align: center;
            color: #fff;
            background-color: #007bff;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            margin-top: 15px;
            text-decoration: none;
        }

        .btn-success {
            background-color: #28a745;
        }

        .btn:hover {
            background-color: #0056b3;
        }

        .btn-success:hover {
            background-color: #218838;
        }

        .card {
            background: #f8f9fa;
            border: 1px solid #e3e3e3;
            border-radius: 10px;
            margin-top: 20px;
            padding: 15px;
        }

        .card-body {
            padding: 15px;
        }

        .card-title {
            font-size: 1.5rem;
            color: #333;
            margin-bottom: 10px;
        }

        .card-text {
            font-size: 1rem;
            color: #555;
        }

        ul {
            list-style-type: disc;
            padding-left: 20px;
        }

        ul li {
            font-size: 1rem;
            color: #555;
            margin-bottom: 5px;
        }
    </style>
</head>
<body>
    <div class="container mt-5">
        <h1 class="text-center mb-4">Study Material Generator</h1>
        <form method="POST" class="mb-5">
            <div class="form-group">
                <label for="topic">Enter the Topic:</label>
                <input type="text" class="form-control" id="topic" name="topic" value="{{ topic }}" placeholder="Enter a topic...">
            </div>
            <button type="submit" class="btn btn-primary btn-block">Generate Study Material</button>
        </form>

        <div class="card">
            <div class="card-body">
                <h4 class="card-title">Summary</h4>
                <p class="card-text">{{ study_material['Summary'] }}</p>
            </div>
        </div>

        <div class="card mt-4">
            <div class="card-body">
                <h4 class="card-title">Key Points</h4>
                <ul>
                    {% for point in study_material['Key Points'] %}
                        <li>{{ point }}</li>
                    {% endfor %}
                </ul>
            </div>
        </div>

        <div class="card mt-4">
            <div class="card-body">
                <h4 class="card-title">Questions</h4>
                <ul>
                    {% for question in study_material['Questions'] %}
                        <li>{{ question }}</li>
                    {% endfor %}
                </ul>
            </div>
        </div>

        <form method="POST" action="/generate_video" onsubmit="handleSubmit(this)">
          <input type="hidden" name="summary" value="{{ study_material['Summary'] }}">
          <input type="hidden" name="key_points" value="{{ study_material['Key Points']|join(',') }}">
          <button type="submit" class="btn btn-success mt-4" id="generateVideoBtn">
              Generate Video Lecture
          </button>
      </form>

      <script>
          function handleSubmit(form) {
              const button = document.getElementById('generateVideoBtn');
              button.disabled = true;
              button.innerHTML = 'Generating... <span class="spinner-border spinner-border-sm"></span>';
          }
      </script>
        <form method="POST" action="/download-pdf">
            <input type="hidden" name="summary" value="{{ study_material['Summary'] }}">
            <input type="hidden" name="key_points" value="{{ study_material['Key Points']|join(',') }}">
            <button type="submit" class="btn btn-success mt-4">Download PDF</button>
        </form>
        <form method="POST" action="/download-ppt">
            <input type="hidden" name="summary" value="{{ study_material['Summary'] }}">
            <input type="hidden" name="key_points" value="{{ study_material['Key Points']|join(',') }}">
            <button type="submit" disable="true" style="pointer-events: none;" class="btn btn-success mt-4">Download PPT (Comming Soon.....)</button>
        </form>
    </div>
</body>
</html>



        ''', study_material=study_material, topic=topic)

    return render_template_string('''
        <!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Study Material Generator</title>
    <style>
        /* General reset */
        body {
            margin: 0;
            padding: 0;
            font-family: 'Arial', sans-serif;
            background: #f9f9f9;
        }

        .container {
            max-width: 600px;
            margin: 50px auto;
            padding: 20px;
            background: #fff;
            border-radius: 10px;
            box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1);
        }

        h1 {
            font-size: 2rem;
            color: #333;
            margin-bottom: 20px;
        }

        .form-group label {
            display: block;
            font-size: 1rem;
            font-weight: 600;
            margin-bottom: 5px;
        }

        .form-group input {
            width: 100%;
            padding: 10px;
            font-size: 1rem;
            border: 1px solid #ddd;
            border-radius: 5px;
            outline: none;
            transition: border 0.3s;
        }

        .form-group input:focus {
            border-color: #007bff;
        }

        button {
            width: 100%;
            padding: 10px;
            font-size: 1rem;
            color: #fff;
            background-color: #007bff;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            transition: background-color 0.3s;
        }

        button:hover {
            background-color: #0056b3;
        }
    </style>
</head>
<body>
    <div class="container">
        <h1 class="text-center">Study Material Generator</h1>
        <form method="POST">
            <div class="form-group">
                <label for="topic">Enter the Topic:</label>
                <input type="text" class="form-control" id="topic" name="topic" placeholder="Enter a topic...">
            </div>
            <button type="submit">Generate Study Material</button>
        </form>
    </div>
</body>
</html>

    ''')

@app.route('/generate_video', methods=['POST'])
def generate_video():
    summary = request.form['summary']
    key_points = request.form['key_points'].split(',')

    video_file = generate_video_from_text_and_images(summary, key_points)

    return send_file(video_file, as_attachment=True)


@app.route('/download-pdf', methods=['POST'])
def download_pdf():
    summary = request.form['summary']
    key_points = request.form['key_points'].split(',')

    pdf_content = f'''
    <h1>Study Material</h1>
    <h2>Summary</h2>
    <p>{summary}</p>
    <h2>Key Points</h2>
    <ul>
    {''.join(f'<li>{point}</li>' for point in key_points)}
    </ul>
    '''

    pdf = pdfkit.from_string(pdf_content, False)

    return send_file(
        io.BytesIO(pdf),
        as_attachment=True,
        download_name='study_material.pdf',
        mimetype='application/pdf'
    )



@app.route('/download-ppt', methods=['POST'])
def download_ppt():
    try:
        summary = request.form['summary']
        key_points = request.form['key_points'].split(',')

        # Create a PowerPoint Presentation
        ppt = Presentation()

        # Add Title Slide
        slide_layout = ppt.slide_layouts[0]  # Title Slide Layout
        slide = ppt.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Study Material"
        subtitle.text = "Generated from API"

        # Add Summary Slide
        if summary:
            slide_layout = ppt.slide_layouts[5]  # Title and Content Layout
            slide = ppt.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Summary"
            content.text = summary

        # Add Key Points Slide
        if key_points:
            slide_layout = ppt.slide_layouts[5]  # Title and Content Layout
            slide = ppt.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Key Points"
            content.text = "\n".join(key_points)  # Join list with new lines

        # Save the PPT to a BytesIO stream
        ppt_stream = io.BytesIO()
        ppt.save(ppt_stream)
        ppt_stream.seek(0)

        # Send the generated PPT as a downloadable file
        return send_file(
            ppt_stream,
            as_attachment=True,
            download_name="study_material.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except json.JSONDecodeError as e:
        return {"error": "Invalid JSON format", "message": str(e)}, 400




# Run app==============================================
if __name__ == '__main__':
    # public_url = ngrok.connect(5000)
    print(f" * Tunnel URL: {public_url}")
    app.run(debug=False, use_reloader=False)
